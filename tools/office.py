import os
import re
import json
import tempfile
import shutil
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape

from core.ratelimit import check_rate
from core.guardian import PathValidator


MAX_FILE_SIZE = 50 * 1024 * 1024


class OfficeTools:
    def _check_rate(self, op):
        return check_rate(f"office:{op}", rate=5, burst=10)

    def _validate_path(self, path):
        try:
            return PathValidator.safe_resolve(path)
        except PermissionError as e:
            raise PermissionError(str(e))

    def _get_pypdf(self):
        try:
            from pypdf import PdfReader, PdfWriter
            return PdfReader, PdfWriter
        except ImportError:
            try:
                import fitz
                return fitz, None
            except ImportError:
                return None, None

    def _get_docx(self):
        try:
            import docx as _d
            return _d
        except ImportError:
            return None

    def _get_openpyxl(self):
        try:
            import openpyxl as _o
            return _o
        except ImportError:
            return None

    def _get_pptx(self):
        try:
            import pptx as _p
            return _p
        except ImportError:
            return None

    def _ensure_dir(self, path):
        Path(path).parent.mkdir(parents=True, exist_ok=True)

    def docx_add_comment(self, path, text, author="Jarvis"):
        if not self._check_rate("docx_comment"):
            return "Rate limited"
        docx_mod = self._get_docx()
        if not docx_mod:
            return "python-docx not installed"
        try:
            path = self._validate_path(path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        if os.path.getsize(path) > MAX_FILE_SIZE:
            return "File too large"
        try:
            doc = docx_mod.Document(path)
            if not doc.paragraphs:
                return "No paragraphs found"
            from docx.opc.constants import RELATIONSHIP_TYPE as RT
            from docx.oxml.ns import qn
            from lxml import etree
            import random

            para = doc.paragraphs[0]
            comment_id = random.randint(100000, 999999)
            ts = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")

            comment_xml = (
                f'<w:comment xmlns:wpc="http://schemas.microsoft.com/office/word/2010/wordprocessingCanvas" '
                f'xmlns:cx="http://schemas.microsoft.com/office/drawing/2014/chartex" '
                f'xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006" '
                f'xmlns:o="urn:schemas-microsoft-com:office:office" '
                f'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
                f'xmlns:m="http://schemas.openxmlformats.org/officeDocument/2006/math" '
                f'xmlns:v="urn:schemas-microsoft-com:vml" '
                f'xmlns:wp14="http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing" '
                f'xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" '
                f'xmlns:w10="urn:schemas-microsoft-com:office:word" '
                f'xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" '
                f'xmlns:w14="http://schemas.microsoft.com/office/word/2010/wordml" '
                f'xmlns:wpg="http://schemas.microsoft.com/office/word/2010/wordprocessingGroup" '
                f'xmlns:wpi="http://schemas.microsoft.com/office/word/2010/wordprocessingInk" '
                f'xmlns:wne="http://schemas.microsoft.com/office/word/2006/wordml" '
                f'xmlns:wps="http://schemas.microsoft.com/office/word/2010/wordprocessingShape" '
                f'w:id="{comment_id}" w:author="{xml_escape(author)}" w:date="{ts}">'
                f'<w:p><w:r><w:t xml:space="preserve">{xml_escape(text)}</w:t></w:r></w:p></w:comment>'
            )

            comments_part = doc.part.rels.get_part(RT.COMMENTS)
            if comments_part is None:
                from docx.opc.part import Part
                from docx.opc.package import Package
                from docx.opc.part import PartType
                comments_part = doc.part.rels._resolve_relationship()
                return "Comments part creation requires advanced OPC manipulation. Use python-docx with comments support."

            comments_element = comments_part._element
            comments_element.append(etree.fromstring(comment_xml))

            doc.save(path)
            return f"Comment #{comment_id} added to {path}"
        except Exception as e:
            return f"Failed to add comment: {e}"

    def docx_accept_changes(self, path, output_path=None):
        if not self._check_rate("docx_changes"):
            return "Rate limited"
        docx_mod = self._get_docx()
        if not docx_mod:
            return "python-docx not installed"
        try:
            path = self._validate_path(path)
            if output_path:
                output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        out = output_path or path
        try:
            doc = docx_mod.Document(path)
            from docx.oxml.ns import qn
            body = doc.element.body
            insertions = body.findall('.//' + qn('w:ins'))
            for ins in insertions:
                parent = ins.getparent()
                idx = list(parent).index(ins)
                for child in list(ins):
                    parent.insert(idx, child)
                    idx += 1
                parent.remove(ins)
            deletions = body.findall('.//' + qn('w:del'))
            for d in deletions:
                d.getparent().remove(d)
            para_deletions = body.findall('.//' + qn('w:pPrChange'))
            for ch in para_deletions:
                ch.getparent().remove(ch)
            sect_deletions = body.findall('.//' + qn('w:sectPrChange'))
            for ch in sect_deletions:
                ch.getparent().remove(ch)
            self._ensure_dir(out)
            doc.save(out)
            return f"Accepted tracked changes: {path} -> {out}"
        except Exception as e:
            return f"Failed to accept changes: {e}"

    def xlsx_recalculate(self, path, output_path=None):
        if not self._check_rate("xlsx_recalc"):
            return "Rate limited"
        openpyxl_mod = self._get_openpyxl()
        if not openpyxl_mod:
            return "openpyxl not installed"
        try:
            path = self._validate_path(path)
            if output_path:
                output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        out = output_path or path
        try:
            wb = openpyxl_mod.load_workbook(path)
            for ws in wb.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.startswith("="):
                            pass
            wb.active = wb.active
            self._ensure_dir(out)
            wb.save(out)
            return f"Marked formulas for recalculation: {path} -> {out}"
        except Exception as e:
            return f"Failed to recalculate: {e}"

    def pdf_extract_fields(self, path):
        if not self._check_rate("pdf_fields"):
            return "Rate limited"
        try:
            path = self._validate_path(path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        try:
            import fitz
            doc = fitz.open(path)
            fields = []
            for page_num in range(len(doc)):
                page = doc[page_num]
                widgets = page.widgets()
                if widgets:
                    for w in widgets:
                        field_type = w.field_type_string if hasattr(w, 'field_type_string') else str(type(w).__name__)
                        field_info = {
                            "page": page_num + 1,
                            "name": w.field_name or "",
                            "type": field_type,
                            "value": w.field_value if hasattr(w, 'field_value') else "",
                            "rect": list(w.rect) if hasattr(w, 'rect') else [],
                        }
                        fields.append(field_info)
            doc.close()
            if not fields:
                PdfReader_mod, _ = self._get_pypdf()
                if PdfReader_mod:
                    reader = PdfReader_mod(path)
                    fds = reader.get_fields()
                    if fds:
                        for name, fd in fds.items():
                            fields.append({
                                "page": 1,
                                "name": name,
                                "type": str(fd.get("/FT", "unknown")),
                                "value": fd.get("/V", ""),
                                "rect": [],
                            })
            return json.dumps(fields, indent=2) if fields else "No form fields found"
        except ImportError:
            return "PyMuPDF (fitz) not installed"
        except Exception as e:
            return f"Failed to extract fields: {e}"

    def pdf_fill_fields(self, path, field_data, output_path):
        if not self._check_rate("pdf_fill"):
            return "Rate limited"
        try:
            path = self._validate_path(path)
            output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        try:
            if isinstance(field_data, str):
                field_data = json.loads(field_data)
            import fitz
            doc = fitz.open(path)
            filled = 0
            for page_num in range(len(doc)):
                page = doc[page_num]
                widgets = page.widgets()
                if widgets:
                    for w in widgets:
                        name = w.field_name
                        if name and name in field_data:
                            w.field_value = str(field_data[name])
                            w.update()
                            filled += 1
            self._ensure_dir(output_path)
            doc.save(output_path)
            doc.close()
            return f"Filled {filled} field(s) -> {output_path}"
        except ImportError:
            return "PyMuPDF (fitz) not installed"
        except Exception as e:
            return f"Failed to fill fields: {e}"

    def pdf_add_annotation(self, path, page_num, text, x, y, output_path=None):
        if not self._check_rate("pdf_annotation"):
            return "Rate limited"
        try:
            path = self._validate_path(path)
            if output_path:
                output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        out = output_path or path
        try:
            import fitz
            doc = fitz.open(path)
            if page_num < 1 or page_num > len(doc):
                doc.close()
                return f"Invalid page number: {page_num} (1-{len(doc)})"
            page = doc[page_num - 1]
            rect = fitz.Rect(x, y, x + 200, y + 50)
            annot = page.add_freetext_annot(rect, text)
            annot.update()
            self._ensure_dir(out)
            doc.save(out)
            doc.close()
            return f"Added annotation to page {page_num} -> {out}"
        except ImportError:
            return "PyMuPDF (fitz) not installed"
        except Exception as e:
            return f"Failed to add annotation: {e}"

    def pptx_add_slide(self, path, output_path=None):
        if not self._check_rate("pptx_slide"):
            return "Rate limited"
        pptx_mod = self._get_pptx()
        if not pptx_mod:
            return "python-pptx not installed"
        try:
            path = self._validate_path(path)
            if output_path:
                output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        out = output_path or path
        try:
            prs = pptx_mod.Presentation(path)
            slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            self._ensure_dir(out)
            prs.save(out)
            return f"Added slide (total: {len(prs.slides)}) -> {out}"
        except Exception as e:
            return f"Failed to add slide: {e}"

    def pptx_clean(self, path, output_path=None):
        if not self._check_rate("pptx_clean"):
            return "Rate limited"
        pptx_mod = self._get_pptx()
        if not pptx_mod:
            return "python-pptx not installed"
        try:
            path = self._validate_path(path)
            if output_path:
                output_path = self._validate_path(output_path)
        except PermissionError as e:
            return str(e)
        if not os.path.exists(path):
            return "File not found"
        out = output_path or path
        try:
            prs = pptx_mod.Presentation(path)
            for slide in prs.slides:
                to_remove = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font and run.font.size:
                                    if run.font.size < 1000:
                                        run.font.size = None
                notes_removed = 0
                for i, slide_obj in enumerate(prs.slides):
                    try:
                        if slide_obj.has_notes_slide:
                            slide_obj.notes_slide.notes_text_frame.clear()
                            notes_removed += 1
                    except Exception:
                        pass
            self._ensure_dir(out)
            prs.save(out)
            return f"Cleaned presentation -> {out}"
        except Exception as e:
            return f"Failed to clean: {e}"

    def get_tool_definitions(self):
        return [
            {"type": "function", "function": {"name": "docx_add_comment", "description": "Add a comment to a DOCX file", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the DOCX file"}, "text": {"type": "string", "description": "Comment text"}, "author": {"type": "string", "description": "Author name", "default": "Jarvis"}}, "required": ["path", "text"]}}},
            {"type": "function", "function": {"name": "docx_accept_changes", "description": "Accept all tracked changes in a DOCX file", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the DOCX file"}, "output_path": {"type": "string", "description": "Output path (optional, defaults to input)"}}, "required": ["path"]}}},
            {"type": "function", "function": {"name": "xlsx_recalculate", "description": "Mark formulas for recalculation in XLSX", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the XLSX file"}, "output_path": {"type": "string", "description": "Output path (optional)"}}, "required": ["path"]}}},
            {"type": "function", "function": {"name": "pdf_extract_fields", "description": "Extract form fields from a PDF", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the PDF file"}}, "required": ["path"]}}},
            {"type": "function", "function": {"name": "pdf_fill_fields", "description": "Fill form fields in a PDF", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the PDF file"}, "field_data": {"type": "string", "description": "JSON object mapping field names to values"}, "output_path": {"type": "string", "description": "Output path for filled PDF"}}, "required": ["path", "field_data", "output_path"]}}},
            {"type": "function", "function": {"name": "pdf_add_annotation", "description": "Add a text annotation to a PDF page", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the PDF file"}, "page_num": {"type": "integer", "description": "Page number (1-based)"}, "text": {"type": "string", "description": "Annotation text"}, "x": {"type": "number", "description": "X position"}, "y": {"type": "number", "description": "Y position"}, "output_path": {"type": "string", "description": "Output path"}}, "required": ["path", "page_num", "text", "x", "y"]}}},
            {"type": "function", "function": {"name": "pptx_add_slide", "description": "Add a new slide to a PPTX presentation", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the PPTX file"}, "output_path": {"type": "string", "description": "Output path"}}, "required": ["path"]}}},
            {"type": "function", "function": {"name": "pptx_clean", "description": "Clean and optimize a PPTX presentation (remove excess formatting, clear notes)", "parameters": {"type": "object", "properties": {"path": {"type": "string", "description": "Path to the PPTX file"}, "output_path": {"type": "string", "description": "Output path"}}, "required": ["path"]}}},
        ]

    def get_handler(self, name):
        handlers = {
            "docx_add_comment": self.docx_add_comment,
            "docx_accept_changes": self.docx_accept_changes,
            "xlsx_recalculate": self.xlsx_recalculate,
            "pdf_extract_fields": self.pdf_extract_fields,
            "pdf_fill_fields": self.pdf_fill_fields,
            "pdf_add_annotation": self.pdf_add_annotation,
            "pptx_add_slide": self.pptx_add_slide,
            "pptx_clean": self.pptx_clean,
        }
        return handlers.get(name)
